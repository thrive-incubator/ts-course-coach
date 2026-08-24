"""Inline pedagogy coaching for course-proposal fields and modules."""

from __future__ import annotations

import io
import json
import logging
import re
from typing import Any

from fastapi import APIRouter, File, Form, HTTPException, UploadFile
from google.genai import types
from pydantic import BaseModel, Field

from app.core.llm import generate_json

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/coach", tags=["coach"])


class CoachRequest(BaseModel):
    section: str
    field: str
    current_value: str = ""
    course_context: dict[str, Any] = {}


class CoachResponse(BaseModel):
    suggestion: str = Field(description="ONE focused suggestion, at most 90 words")
    examples: list[str] = Field(default_factory=list, description="2-3 short examples, at most 30 words each")
    tone: str = Field(default="encouraging", description="encouraging | challenging | celebratory")


class SiblingModule(BaseModel):
    index: int
    module_name: str = ""
    contact_hours: str = ""
    is_current: bool = False


class ModuleReviewRequest(BaseModel):
    module: dict[str, Any]
    course_essential_question: str = ""
    course_learning_objectives: str = ""
    sibling_modules: list[SiblingModule] = []
    course_context: dict[str, Any] = {}


class ModuleReviewResponse(BaseModel):
    strengths: list[str] = Field(
        default_factory=list,
        description="2-4 specific things that are working — cite exact language when possible. At most 60 words each.",
    )
    gaps: list[str] = Field(
        default_factory=list,
        description="2-4 specific things missing or weak — be direct, not preachy. At most 60 words each.",
    )
    suggestions: list[str] = Field(
        default_factory=list,
        description="2-4 concrete moves the faculty could make this session. At most 60 words each.",
    )
    bloom_diagnosis: str = Field(
        default="",
        description="1-2 sentences on where objectives cluster on Bloom's and one specific rewrite showing sharper verbs",
    )
    interactive_ideas: list[str] = Field(
        default_factory=list,
        description="2-3 specific interactive-feature ideas grounded in THIS module's content. At most 60 words each.",
    )


class MaterialReviewLLM(BaseModel):
    summary: str = Field(description="1-2 sentence read of what this material is doing right now (not what it should do)")
    strengths: list[str] = Field(
        default_factory=list,
        description="2-4 things the material does well — cite specific slides/pages/sections when possible. At most 60 words each.",
    )
    improvements: list[str] = Field(
        default_factory=list,
        description="2-4 concrete pedagogy improvements — each specific enough to act on. At most 60 words each.",
    )
    bloom_diagnosis: str = Field(
        default="",
        description="1-2 sentences on where the material lives on Bloom's and how to push it up one level",
    )
    engagement_ideas: list[str] = Field(
        default_factory=list,
        description="2-3 specific engagement moves that would work with THIS content. At most 60 words each.",
    )


class MaterialReviewResponse(MaterialReviewLLM):
    extracted_chars: int = 0


FIELD_PROMPTS: dict[str, str] = {
    "course_description": (
        "You are coaching a faculty member writing a course description for Thrive Academy. "
        "Read what they've written and offer one focused suggestion that improves clarity of "
        "value-to-learner AND that a prospective enrollee could scan in 15 seconds. "
        "If empty, propose a strong 2-3 sentence opener based on the course name and audience. "
        "Then offer 2 short example sentences they could remix."
    ),
    "needs_statement": (
        "You are coaching a faculty member on their Needs Statement. The best needs statements "
        "name a specific practitioner or systemic gap and cite (or gesture toward) evidence. "
        "Weak ones describe the topic instead of demand. Read the draft. If it's a topic "
        "description rather than a demand claim, say so plainly and rewrite one sentence to "
        "show the shape. If it's already strong, name one sharpening move. "
        "Give 2 example openers grounded in the course topic."
    ),
    "evidence_of_demand": (
        "Coach the faculty on Evidence of Demand. Strong evidence cites: prior waitlists, survey "
        "data, workforce reports, employer requests, or CEU-requirement changes. If the draft "
        "is vague, ask (in your suggestion) which of those sources they could add. "
        "Suggest 2 concrete evidence types that would fit this course topic."
    ),
    "competitive_landscape": (
        "Coach the faculty on Competitive Landscape. The trap is 'no competitors' — every "
        "learner has alternatives (a YouTube series, a $5K CEU program, a webinar). Push them "
        "to name 2-3 real alternatives and their positioning. Then help them articulate the "
        "distinct wedge Thrive Academy offers. Suggest 2 sample alternative-plus-differentiator lines."
    ),
    "recruitment_and_marketing": (
        "Coach the faculty on Recruitment & Marketing. Focus on: (1) where the audience already "
        "gathers (listservs, associations, LinkedIn groups, conferences), (2) the one line that "
        "would make someone in that audience stop scrolling, (3) proof points that unlock trust. "
        "Suggest 2 candidate one-liners tuned to the intended audience."
    ),
    "essential_question": (
        "Coach the faculty on their COURSE-LEVEL Essential Question. A strong essential "
        "question is: (a) open, not answerable with a fact; (b) provocative — it invites "
        "wrestling; (c) durable — worth returning to across many modules; (d) rooted in what "
        "makes the practitioner community actually change. Read the draft. If it's really a "
        "learning objective in disguise (starts with 'Understand…' or 'Know…'), name it and "
        "rewrite as a genuine question. Give 3 example essential questions tuned to this "
        "course's audience and topic, at varying levels of provocation."
    ),
    "learning_objectives": (
        "Coach the faculty on Learning Objectives using Bloom's Revised Taxonomy. Objectives "
        "should start with an observable verb at the appropriate cognitive level "
        "(remember/understand/apply/analyze/evaluate/create). Read their draft. If any objective "
        "uses vague verbs ('know', 'understand', 'appreciate'), name it and rewrite one with a "
        "sharper Bloom's verb. Suggest 3 example objectives at varying Bloom levels for this course."
    ),
    "course_structure": (
        "Coach the faculty on Course Structure. Strong structures show how the learning journey "
        "arcs (foundational -> applied -> synthesized), match the format (async/sync/live), and "
        "surface where practice + feedback loops live. Give one specific structural suggestion "
        "based on the format, duration and contact hours."
    ),
    "assessment_methods": (
        "Coach the faculty on Grading Scheme & Assessment. Innovative pedagogy favors: authentic "
        "assessments (real-world artifacts), formative feedback loops, peer review, and rubrics "
        "shared upfront. Lecture-then-quiz is the weakest form. Suggest 2 assessment "
        "methods that would fit this course type and be more engaging than a final quiz."
    ),
    "student_support": (
        "Coach the faculty on Student Support. Faculty often under-describe support. Prompt them "
        "to name: office hours cadence, cohort community mechanisms (Slack? forum?), 1:1 access, "
        "and any TA/coach layer. Suggest 2 support elements that match the course format."
    ),
    "cqi": (
        "Coach the faculty on CQI & Staying Current. Ask them to name: mid-course pulse survey, "
        "end-of-cohort NPS, an update cadence for materials (quarterly? per-cohort?), and how "
        "they'll incorporate emergent research. Suggest 2 concrete CQI mechanisms."
    ),
    "financial_overview": (
        "Coach the faculty on Financial Overview. Ask: what's the tuition, cohort size, break-even "
        "cohort size, major cost drivers (faculty time, platform, marketing), and 3-cohort revenue "
        "projection? Use the tuition and cohort size from the course context if they are given "
        "there. If any are missing, name them. If the plan feels thin, suggest one "
        "sensitivity to test (e.g., 'what if cohort fills at 70%?')."
    ),
}

_GENERIC_PROMPT = (
    "You are coaching a faculty member on their course proposal. Read the draft field content, "
    "consider the course context, and offer one focused, actionable suggestion that improves "
    "either pedagogical strength or marketability. Then offer 2 short examples they could remix."
)


def _system_instruction() -> str:
    return (
        "You are the Course Coach for Thrive Academy — an experienced learning designer + "
        "instructional marketer. Your voice is: warm, concise, concrete. You NEVER lecture. "
        "You give ONE main suggestion per response (at most 90 words), then 2-3 short examples "
        "(at most 30 words each — the field prompt says how many). You reference "
        "innovative-pedagogy principles (Bloom's, authentic assessment, active learning, "
        "cohort community) without name-dropping them mechanically. When marketing, you focus "
        "on audience-first framing, not feature lists.\n\n"
        "If the draft is empty, do not critique the absence — propose a strong starter built "
        "from the course context, then the examples. If the draft is already strong, say so in "
        "one line and name one sharpening move; never invent a problem.\n\n"
        "Example of the voice, for a course-level essential question drafted as "
        "'Understand reflective supervision':\n"
        '  suggestion: "That\'s a learning objective wearing a question\'s clothes — it can be '
        "answered with a definition. An essential question should keep a home visitor arguing "
        "with herself in week ten. Try aiming it at the moment reflective practice is hardest: "
        'when the family is in crisis and slowing down feels irresponsible."\n'
        '  examples: ["What does it take to stay curious about a parent when everything in you '
        'wants to fix?", "When is slowing down the most responsible thing a home visitor can do?"]\n'
        '  tone: "challenging"\n\n'
        "Return valid JSON only."
    )


def _module_system_instruction() -> str:
    return (
        "You are the Course Coach for Thrive Academy — a senior learning designer who reviews "
        "individual course modules with a pedagogy-first lens. You look at essential questions, "
        "learning objectives (Bloom's alignment), critical information, engagement, and "
        "interactive features together, and at where the module sits in the course sequence. "
        "You are warm, concrete, and specific. You never lecture. "
        "You never say 'consider adding…' — you show an example. Each bullet is at most 60 "
        "words. If a section of the module is empty, propose a starter for it instead of "
        "criticising the gap. Return valid JSON only."
    )


def _material_system_instruction() -> str:
    return (
        "You are the Course Coach for Thrive Academy reviewing a lesson-plan or slide deck a "
        "faculty member uploaded for a specific module. You read the material and evaluate "
        "how it lands as a learning experience: Are learning objectives clear? Do activities go "
        "beyond passive consumption? Where does Bloom's stall at 'remember/understand'? What "
        "engagement moves could deepen learning? You are warm, concrete, non-preachy. Each "
        "bullet is at most 60 words and points to a specific slide, page or section when it "
        "can. Return valid JSON only."
    )


def _prompt_for(field: str) -> str:
    prompt = FIELD_PROMPTS.get(field)
    if prompt is None:
        logger.warning("coach: no field prompt for %r — using generic prompt", field)
        return _GENERIC_PROMPT
    return prompt


def _context_lines(ctx: dict[str, Any] | None) -> str:
    return "\n".join(f"  {k}: {v}" for k, v in (ctx or {}).items() if v)


@router.post("", response_model=CoachResponse)
async def coach(req: CoachRequest) -> CoachResponse:
    """Return an inline coaching suggestion for a single proposal field."""
    field_prompt = _prompt_for(req.field)
    user_text = (
        f"Field: {req.field} (in section: {req.section})\n\n"
        f"Course context so far:\n{_context_lines(req.course_context) or '  (none yet)'}\n\n"
        f"Current draft of this field:\n---\n{req.current_value or '(empty — no draft yet)'}\n---\n\n"
        f"{field_prompt}\n\n"
        'Reply with strict JSON: {"suggestion": "...", "examples": ["...", "..."], '
        '"tone": "encouraging|challenging|celebratory"}'
    )
    result = await generate_json(
        endpoint=f"coach/{req.field}",
        system=_system_instruction(),
        contents=user_text,
        schema=CoachResponse,
        temperature=0.7,
        # Low thinking: ~2s instead of ~7s on gemini-3.x flash, same quality for
        # a single-field nudge. Deeper reviews below keep default thinking.
        thinking_level="low",
    )
    result.suggestion = result.suggestion[:2000]
    result.examples = [e[:500] for e in result.examples][:5]
    if result.tone not in {"encouraging", "challenging", "celebratory"}:
        result.tone = "encouraging"
    return result


def _module_summary(mod: dict[str, Any]) -> str:
    objectives = mod.get("objectives") or []
    obj_lines = "\n".join(
        f"    - [{(o.get('bloom') or '?')}] {o.get('text') or '(blank)'}"
        for o in objectives
    ) or "    (none written yet)"
    features = ", ".join(mod.get("interactive_features") or []) or "(none selected)"
    return (
        f"Module name: {mod.get('module_name') or '(untitled)'}\n"
        f"Contact hours: {mod.get('contact_hours') or '?'}\n"
        f"Format: {mod.get('format') or '?'}\n"
        f"Module essential question: {mod.get('essential_question') or '(none written)'}\n"
        f"Learning objectives (with Bloom's level):\n{obj_lines}\n"
        f"Critical information: {mod.get('critical_information') or '(none)'}\n"
        f"Engagement opportunities: {mod.get('engagement_opportunities') or '(none)'}\n"
        f"Interactive features selected: {features}\n"
        f"Notes on interactive features: {mod.get('interactive_features_notes') or '(none)'}\n"
        f"Assignments: {mod.get('assignments') or '(none)'}\n"
    )


def _sibling_lines(siblings: list[SiblingModule]) -> str:
    if not siblings:
        return "  (only this module so far)"
    return "\n".join(
        f"  {s.index}. {s.module_name or '(untitled)'}"
        f"{f' — {s.contact_hours} hrs' if s.contact_hours else ''}"
        f"{'   <== this module' if s.is_current else ''}"
        for s in siblings
    )


def _course_frame(
    course_context: dict[str, Any] | None,
    course_essential_question: str,
    course_learning_objectives: str = "",
    siblings: list[SiblingModule] | None = None,
) -> str:
    return (
        f"COURSE CONTEXT:\n{_context_lines(course_context) or '  (none)'}\n\n"
        f"COURSE ESSENTIAL QUESTION: {course_essential_question or '(not yet written)'}\n\n"
        f"COURSE-LEVEL LEARNING OBJECTIVES:\n{course_learning_objectives or '  (not yet written)'}\n\n"
        f"MODULE SEQUENCE:\n{_sibling_lines(siblings or [])}\n\n"
    )


@router.post("/module", response_model=ModuleReviewResponse)
async def coach_module(req: ModuleReviewRequest) -> ModuleReviewResponse:
    """Holistic review of a single module in context of the course."""
    user_text = (
        _course_frame(
            req.course_context,
            req.course_essential_question,
            req.course_learning_objectives,
            req.sibling_modules,
        )
        + f"MODULE DRAFT:\n{_module_summary(req.module)}\n\n"
        "Review this module holistically. Focus on:\n"
        "  1) Does the module essential question ladder up to the course essential question?\n"
        "  2) Do the learning objectives use observable Bloom's verbs? Are they at appropriate "
        "levels for where this module sits in the sequence (earlier modules lower, later higher)?\n"
        "  3) Is 'critical information' concept-rich or just topic-listy?\n"
        "  4) Do engagement opportunities create wrestle, or is this passive consumption?\n"
        "  5) Do the selected interactive features (simulation, case study, generative chat, etc.) "
        "match the objectives?\n"
        "  6) Does it overlap with, or leave a gap against, the neighbouring modules?\n\n"
        "Return strict JSON with strengths, gaps, suggestions, bloom_diagnosis, interactive_ideas."
    )
    result = await generate_json(
        endpoint="coach/module",
        system=_module_system_instruction(),
        contents=user_text,
        schema=ModuleReviewResponse,
        temperature=0.6,
    )
    return _cap_lists(result)


def _cap_lists(result: Any, cap: int = 500, n: int = 6) -> Any:
    for name, value in list(result.__dict__.items()):
        if isinstance(value, list):
            setattr(result, name, [str(x)[:cap] for x in value][:n])
        elif isinstance(value, str) and name != "summary":
            setattr(result, name, value[:1200])
    return result


# ---- File extraction ---------------------------------------------------------

MAX_UPLOAD_BYTES = 15 * 1024 * 1024  # 15MB
# Gemini 3.x flash has a very large context; this cap only guards against pathological files.
MAX_TEXT_CHARS = 200_000


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:  # noqa: BLE001
            continue
        if text.strip():
            parts.append(f"--- Page {i} ---\n{text}")
    return "\n\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        chunks: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        chunks.append(line)
        try:
            notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        except Exception:  # noqa: BLE001
            notes = ""
        body = "\n".join(chunks)
        if notes:
            body += f"\n\n[Speaker notes]\n{notes}"
        slides.append(f"--- Slide {i} ---\n{body}")
    return "\n\n".join(slides)


def _extract_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def _ext(filename: str) -> str:
    return filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


def _extract_text(filename: str, data: bytes) -> str:
    ext = _ext(filename)
    try:
        if ext == "pdf":
            return _extract_pdf(data)
        if ext == "pptx":
            return _extract_pptx(data)
        if ext == "docx":
            return _extract_docx(data)
        if ext in {"txt", "md", "markdown"}:
            return data.decode("utf-8", errors="replace")
    except Exception as exc:  # noqa: BLE001
        logger.exception("text extraction failed for %s", filename)
        if "decrypt" in str(exc).lower() or "encrypted" in str(exc).lower():
            raise HTTPException(
                status_code=415, detail=f"{filename} is password-protected. Please remove the password and upload again."
            ) from exc
        raise HTTPException(status_code=415, detail=f"Could not read {filename}: {exc}") from exc
    raise HTTPException(
        status_code=415, detail=f"Unsupported file type: .{ext or '?'} — please upload a PDF, PPTX, DOCX or TXT."
    )


@router.post("/material", response_model=MaterialReviewResponse)
async def coach_material(
    file: UploadFile = File(...),
    module: str = Form("{}"),
    course_essential_question: str = Form(""),
    course_context: str = Form("{}"),
) -> MaterialReviewResponse:
    """Review an uploaded lesson plan / slide deck and return coach feedback.

    PDFs go to Gemini natively (it reads layout, diagrams and image-only pages);
    other formats are text-extracted first.
    """
    raw = await file.read()
    if len(raw) > MAX_UPLOAD_BYTES:
        raise HTTPException(status_code=413, detail=f"File too large ({len(raw) // (1024*1024)} MB). Max is 15 MB.")

    filename = file.filename or "upload"
    try:
        module_dict = json.loads(module) if module else {}
        ctx_dict = json.loads(course_context) if course_context else {}
    except json.JSONDecodeError as exc:
        raise HTTPException(status_code=400, detail=f"Bad JSON in form field: {exc}") from exc

    review_task = (
        "Review this material as a coach would. Focus on whether it stretches learners "
        "beyond passive consumption, whether it lands the module's essential question, and "
        "where Bloom's is stalling. Return strict JSON with summary, strengths, improvements, "
        "bloom_diagnosis, engagement_ideas."
    )
    frame = _course_frame(ctx_dict, course_essential_question) + (
        f"MODULE THIS MATERIAL IS FOR:\n{_module_summary(module_dict)}\n\n"
        f"UPLOADED MATERIAL: {filename}\n"
    )

    contents: str | list[Any]
    if _ext(filename) == "pdf":
        # Native PDF: no text extraction, so image-only slides and diagrams are visible to the model.
        extracted_chars = 0
        contents = [
            types.Part.from_text(text=frame + "The material is attached as a PDF.\n\n" + review_task),
            types.Part.from_bytes(data=raw, mime_type="application/pdf"),
        ]
    else:
        text = _extract_text(filename, raw)
        text = re.sub(r"\n{3,}", "\n\n", text).strip()
        if not text:
            raise HTTPException(
                status_code=422,
                detail="No readable text found in that file. If it's a scanned document, export it as a PDF and try again.",
            )
        extracted_chars = len(text)
        truncated = text[:MAX_TEXT_CHARS]
        was_truncated = len(text) > MAX_TEXT_CHARS
        contents = (
            frame
            + (f"(text truncated at {MAX_TEXT_CHARS} chars — reviewing what we could read)\n" if was_truncated else "")
            + f"--- Extracted text ---\n{truncated}\n--- End of text ---\n\n"
            + review_task
        )

    result = await generate_json(
        endpoint="coach/material",
        system=_material_system_instruction(),
        contents=contents,
        schema=MaterialReviewLLM,
        temperature=0.6,
    )
    result = _cap_lists(result)
    return MaterialReviewResponse(**result.model_dump(), extracted_chars=extracted_chars)
